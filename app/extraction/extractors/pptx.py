from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.extraction.extractors.base import Extractor


class PptxExtractor(Extractor):
    file_type = "pptx"

    def extract_elements(self, path: str) -> list[dict[str, Any]]:
        prs = Presentation(path)
        elements: list[dict[str, Any]] = []
        for idx, slide in enumerate(prs.slides, start=1):
            items: list[dict[str, Any]] = []
            for shape in slide.shapes:
                _walk_shape(shape, items)
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text and notes_text.strip():
                    items.append({"type": "paragraph", "text": notes_text, "source": "notes"})
            elements.append({"type": "slide", "index": idx, "items": items})
        return elements


def _walk_shape(shape, items: list[dict[str, Any]]) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            _walk_shape(sub, items)
        return
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs) or paragraph.text
            if text.strip():
                items.append({"type": "paragraph", "text": text})
    if shape.has_table:
        rows: list[list[str]] = []
        for row in shape.table.rows:
            rows.append([cell.text for cell in row.cells])
        items.append({"type": "table", "rows": rows})
