import shutil
from pathlib import Path

import pytest

FIXTURES = Path(__file__).parent / "fixtures"


@pytest.fixture(scope="session", autouse=True)
def build_fixtures() -> None:
    FIXTURES.mkdir(parents=True, exist_ok=True)
    _make_docx(FIXTURES / "sample.docx")
    _make_docx(FIXTURES / "sample.docm")
    _make_xlsx(FIXTURES / "sample.xlsx")
    _make_xlsx(FIXTURES / "sample.xlsm")
    _make_pptx(FIXTURES / "sample.pptx")
    _make_pptx(FIXTURES / "sample.pptm")
    _make_pdf(FIXTURES / "sample.pdf")
    _make_two_column_pdf(FIXTURES / "sample_two_col.pdf")
    (FIXTURES / "sample.txt").write_text("Plain text file\nWith two lines.", encoding="utf-8")
    (FIXTURES / "sample.md").write_text("# Heading\n\nBody paragraph.", encoding="utf-8")
    _make_legacy_via_soffice()


def _make_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_paragraph("Hello docx world")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "A1"
    table.rows[0].cells[1].text = "B1"
    table.rows[1].cells[0].text = "A2"
    table.rows[1].cells[1].text = "B2"
    doc.add_paragraph("After table")
    doc.save(str(path))


def _make_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["name", "value"])
    ws.append(["alpha", 1])
    ws.append(["beta", 2])
    wb.save(str(path))


def _make_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Slide title"
    slide.placeholders[1].text = "Slide subtitle"
    prs.save(str(path))


def _make_pdf(path: Path) -> None:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello pdf world")
    doc.save(str(path))
    doc.close()


def _make_two_column_pdf(path: Path) -> None:
    """Two-column prose page. Column order must be preserved after extraction."""
    import fitz

    left = (
        "LEFTSTART. This is the opening paragraph of the left column. "
        "It continues with several sentences that flow downward naturally. "
        "The left column contains a substantial amount of prose content "
        "to simulate a real newspaper or academic paper layout. "
        "LEFTEND marks the last visible word of the left column."
    )
    right = (
        "RIGHTSTART. This is the opening paragraph of the right column. "
        "It should appear AFTER all left-column content in any correct "
        "reading-order extraction. More filler prose follows here to "
        "ensure the two columns are not mistaken for a small table. "
        "RIGHTEND marks the last visible word of the right column."
    )
    doc = fitz.open()
    page = doc.new_page(width=612, height=792)
    page.insert_textbox(fitz.Rect(50, 50, 290, 750), left, fontsize=11)
    page.insert_textbox(fitz.Rect(322, 50, 562, 750), right, fontsize=11)
    doc.save(str(path))
    doc.close()


def _make_legacy_via_soffice() -> None:
    if not shutil.which("soffice"):
        return
    import subprocess

    mapping = {
        "sample.docx": ("doc", "sample.doc"),
        "sample.xlsx": ("xls", "sample.xls"),
        "sample.pptx": ("ppt", "sample.ppt"),
    }
    for src, (fmt, dst) in mapping.items():
        src_path = FIXTURES / src
        dst_path = FIXTURES / dst
        if dst_path.exists():
            continue
        subprocess.run(
            [
                "soffice",
                "--headless",
                "--norestore",
                "--nologo",
                "--nolockcheck",
                "--convert-to",
                fmt,
                "--outdir",
                str(FIXTURES),
                str(src_path),
            ],
            capture_output=True,
            timeout=60,
        )
